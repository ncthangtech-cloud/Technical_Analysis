"""
Streamlit PDF, pptx Question-Answering App (OpenAI >= 1.0.0 API)

Features:
1) Read multiple PDF, pptx files (uploaded via Streamlit)
2) Chunk the text with overlap
3) Create embeddings using OpenAI Embeddings API
4) Use a FAISS to retrieve relevant chunks
5) Ask OpenAI ChatCompletion to answer user's question using retrieved context

Requirements:
pip install streamlit PyPDF2 openai numpy scikit-learn

Environment:
Set OPENAI_API_KEY environment variable before running.

Run:
streamlit run RR-V01.py

"""

import time
from typing import List, Tuple

import streamlit as st
import numpy as np
#from sklearn.metrics.pairwise import cosine_similarity
from PyPDF2 import PdfReader
from openai import OpenAI
from pptx import Presentation  # NEW import
import os
import glob
import faiss
import pickle
from datetime import datetime

# -----------------------------
# Configuration
# -----------------------------
EMBEDDING_MODEL = "text-embedding-3-small"  # OpenAI embedding model
CHAT_MODEL = "gpt-4o"  # change if you prefer another chat model
# Initialize client
client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

CACHE_PATH = "embeddings_store.pkl"

# ---- LOGIN SETUP ----
USER_CREDENTIALS = st.secrets["users"]  # Read from secrets.toml

# Initialize login state
if "authenticated" not in st.session_state:
    st.session_state.authenticated = False
if "username" not in st.session_state:
    st.session_state.username = ""

# Login form (only show if not logged in)
if not st.session_state.authenticated:
    st.subheader("🔐 Please log in to access the app")

    with st.form("login_form"):
        username = st.text_input("Username")
        password = st.text_input("Password", type="password")
        submit = st.form_submit_button("Login")

        if submit:
            if username in USER_CREDENTIALS and USER_CREDENTIALS[username] == password:
                st.session_state.authenticated = True
                st.session_state.username = username
                st.success(f"Welcome, {username}!")
                st.rerun()  # refresh UI
            else:
                st.error("Invalid username or password")

# Show main app if logged in
if st.session_state.authenticated:
    st.sidebar.success(f"Logged in as {st.session_state.username}")
    if st.sidebar.button("Logout"):
        st.session_state.authenticated = False
        st.session_state.username = ""
        st.rerun()

    def load_cache():
        if os.path.exists(CACHE_PATH):
            with open(CACHE_PATH, "rb") as f:
                return pickle.load(f)
        return {"files": {}, "chunks": [], "embeddings": None, "index": None}
    
    def save_cache(store):
        with open(CACHE_PATH, "wb") as f:
            pickle.dump(store, f)
    
    def process_folder(folder_path, chunk_size=1000, chunk_overlap=200):
        # Load existing cache
        store = load_cache()
        processed_files = store["files"]
    
        # Find new/updated files
        files = glob.glob(os.path.join(folder_path, "*.pdf")) + glob.glob(os.path.join(folder_path, "*.pptx"))
        new_chunks, new_embeds = [], []
    
        for file_path in files:
            mtime = os.path.getmtime(file_path)
            key = (os.path.basename(file_path), mtime)
            if file_path in processed_files and processed_files[file_path] == mtime:
                continue  # already processed
    
            # Extract text
            if file_path.lower().endswith(".pdf"):
                text = extract_text_from_pdf(file_path)
            else:
                text = extract_text_from_pptx(file_path)
    
            chunks = chunk_text(text, chunk_size, chunk_overlap)
            embeddings = get_embeddings(chunks)
    
            new_chunks.extend(chunks)
            new_embeds.append(embeddings)
    
            # Update cache metadata
            processed_files[file_path] = mtime
    
        # If new embeddings found → update FAISS
        if new_embeds:
            new_embeds = np.vstack(new_embeds).astype("float32")
    
            if store["index"] is None:
                d = new_embeds.shape[1]
                index = faiss.IndexFlatIP(d)
                index.add(new_embeds)
                store["index"] = index
                store["chunks"] = new_chunks
            else:
                store["index"].add(new_embeds)
                store["chunks"].extend(new_chunks)
    
        # Save updated cache
        save_cache(store)
        return store
    
    # -----------------------------
    # Utilities: PDF reading and chunking
    # -----------------------------
    
    def extract_text_from_pdf(file_path: str) -> str:
        reader = PdfReader(file_path)
        pages = []
        for p in reader.pages:
            try:
                pages.append(p.extract_text() or "")
            except Exception:
                pages.append("")
        return "\n".join(pages)
    
    # -----------------------------
    # Utilities: PPTX reading
    # -----------------------------
    def extract_text_from_pptx(file_path: str) -> str:
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = shape.text.strip()
                    if txt:
                        texts.append(txt)
        return "\n".join(texts)
    
    def chunk_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
        text = text.replace("\n", " ")
        if len(text) <= chunk_size:
            return [text]
    
        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end]
            chunks.append(chunk.strip())
            if end >= len(text):
                break
            start = end - chunk_overlap
        return chunks
    
    # -----------------------------
    # Embeddings and vector store (in-memory)
    # -----------------------------
    
    def get_embeddings(texts: List[str], model: str = EMBEDDING_MODEL) -> np.ndarray:
        embeddings = []
        for t in texts:
            if not t:
                embeddings.append(np.zeros(1536))
                continue
            resp = client.embeddings.create(model=model, input=t)
            emb = np.array(resp.data[0].embedding)
            embeddings.append(emb)
            time.sleep(0.05)
        return np.vstack(embeddings)
    
    FAISS_INDEX_PATH = "faiss_index.bin"
    META_PATH = "embeddings_meta.pkl"  # stores chunks + processed_files metadata
    
    class FAISSVectorStore:
        def __init__(self, index: faiss.Index, chunks: List[str], processed_files: dict):
            """
            index: faiss.Index with normalized vectors (float32)
            chunks: list of chunk strings (order matches vectors in the index)
            processed_files: dict mapping file_path -> mtime (so we can skip processed files)
            """
            self.index = index
            self.chunks = list(chunks)
            self.processed_files = dict(processed_files)
    
        @classmethod
        def build(cls, embeddings: np.ndarray, chunks: List[str], processed_files: dict = None):
            emb = embeddings.astype("float32")
            # normalize rows to unit length for cosine similarity using inner product
            norms = np.linalg.norm(emb, axis=1, keepdims=True)
            norms[norms == 0] = 1.0
            emb = emb / norms
    
            d = emb.shape[1]
            index = faiss.IndexFlatIP(d)
            index.add(emb)
            return cls(index, chunks, processed_files or {})
    
        def add(self, embeddings: np.ndarray, new_chunks: List[str], file_mtimes: dict = None):
            """Add embeddings (numpy) and corresponding chunks to the index and chunks list."""
            emb = embeddings.astype("float32")
            norms = np.linalg.norm(emb, axis=1, keepdims=True)
            norms[norms == 0] = 1.0
            emb = emb / norms
            self.index.add(emb)
            self.chunks.extend(new_chunks)
            if file_mtimes:
                self.processed_files.update(file_mtimes)
    
        def query(self, query_emb: np.ndarray, top_k: int = 5) -> List[Tuple[int, float]]:
            """Return list of (index, score) pairs (score is inner product ~ cosine)."""
            q = np.array(query_emb).astype("float32")
            qn = q / (np.linalg.norm(q) or 1.0)
            D, I = self.index.search(qn.reshape(1, -1), top_k)
            results = []
            for score, idx in zip(D[0], I[0]):
                if idx == -1:
                    continue
                results.append((int(idx), float(score)))
            return results
    
        def save(self, index_path: str = FAISS_INDEX_PATH, meta_path: str = META_PATH):
            faiss.write_index(self.index, index_path)
            meta = {
                "chunks": self.chunks,
                "processed_files": self.processed_files
            }
            with open(meta_path, "wb") as f:
                pickle.dump(meta, f)
    
        @classmethod
        def load(cls, index_path: str = FAISS_INDEX_PATH, meta_path: str = META_PATH):
            if not (os.path.exists(index_path) and os.path.exists(meta_path)):
                return None
            index = faiss.read_index(index_path)
            with open(meta_path, "rb") as f:
                meta = pickle.load(f)
            return cls(index, meta.get("chunks", []), meta.get("processed_files", {}))
    # -----------------------------
    # Answer generation using retrieved context
    # -----------------------------
    
    def answer_question_with_context(question: str, vector_store, chunks, top_k: int = 3):
        """Retrieve top_k chunks, build a prompt, and get a concise answer."""
        try:
            # Embed the question
            q_emb = get_embeddings([question])[0]
    
            # Retrieve top_k most relevant chunks
            results = vector_store.query(q_emb, top_k=top_k)
            retrieved_chunks = [chunks[idx] for idx, _ in results]
            context = "\n\n".join(retrieved_chunks)
    
            # Concise system prompt
            system_prompt = (
                "You are a helpful assistant. "
                "Answer the user's question clearly and as concisely as possible (max 300 words). "
                "Only use the provided context. "
                "If the answer is not in the context, say 'I don't know, can you please write your question more specifically'."
            )
    
            # Call OpenAI API
            response = client.chat.completions.create(
                model="gpt-4o-mini",  # or gpt-4o if available
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": f"Question: {question}\n\nContext:\n{context}"}
                ],
                max_tokens=200,  # hard cap on answer length
                temperature=0.3
            )
    
            return response.choices[0].message.content.strip()
    
        except Exception as e:
            st.error(f"Error while answering: {e}")
            return None
    
    # -----------------------------
    # Streamlit UI
    # -----------------------------
    
    st.set_page_config(page_title="M&E Analysis", layout="wide")
    st.image("vna.png", width=200)
    st.title("M&E Analysis")
    st.markdown(
        "This App allows you to analyse technical data of VNA fleet based on available data using OpenAI. Click "Process data" first, then ask questions."
    )
    
    folder_path = os.path.join(os.path.dirname(__file__), "data")
    
    #chunk_size = st.sidebar.number_input("Chunk size (chars)", min_value=200, max_value=10000, value=5000, step=100)
    #chunk_overlap = st.sidebar.number_input("Chunk overlap (chars)", min_value=0, max_value=2000, value=500, step=50)
    #top_k = st.sidebar.slider("Top K retrieved chunks", min_value=1, max_value=10, value=4)
    chunk_size = 5000
    chunk_overlap = 500
    top_k = 4
    
    if 'vector_store' not in st.session_state:
        st.session_state.vector_store = None
        st.session_state.chunks = None
    
    # put this where your current Process folder logic is
    if st.button("Process data"):
        if not os.path.isdir(folder_path):
            st.error("Invalid folder path")
        else:
            # try to load existing FAISS store (index + meta)
            faiss_store = FAISSVectorStore.load()  # None if not present
    
            # discover files
            files = glob.glob(os.path.join(folder_path, "*.pdf")) + glob.glob(os.path.join(folder_path, "*.pptx"))
            if not files:
                st.warning("No PDF or PPTX files found in this folder.")
            else:
                # track newly created chunks & embeddings (for adding to index)
                all_new_chunks = []
                all_new_embeds = []
                new_file_mtimes = {}
    
                for file_path in files:
                    mtime = os.path.getmtime(file_path)
                    already_processed = faiss_store and (file_path in faiss_store.processed_files and faiss_store.processed_files[file_path] == mtime)
                    if already_processed:
                        #st.write(f"Skipping already processed: {os.path.basename(file_path)}")
                        continue
    
                    st.write(f"Processing new/updated file: {os.path.basename(file_path)}")
                    if file_path.lower().endswith(".pdf"):
                        text = extract_text_from_pdf(file_path)
                    else:
                        text = extract_text_from_pptx(file_path)
    
                    chunks = chunk_text(text, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
                    if not chunks:
                        continue
    
                    # create embeddings for these chunks
                    embeds = get_embeddings(chunks).astype("float32")
    
                    all_new_chunks.extend(chunks)
                    all_new_embeds.append(embeds)
                    new_file_mtimes[file_path] = mtime
    
                if all_new_embeds:
                    new_embeds = np.vstack(all_new_embeds).astype("float32")
    
                    if faiss_store is None:
                        # first time: build a new FAISS store
                        faiss_store = FAISSVectorStore.build(new_embeds, all_new_chunks, processed_files=new_file_mtimes)
                    else:
                        # add to existing store
                        faiss_store.add(new_embeds, all_new_chunks, file_mtimes=new_file_mtimes)
    
                    # persist to disk
                    try:
                        faiss_store.save()
                    except Exception as e:
                        st.warning(f"Failed to save persistent index: {e}")
    
                # put the wrapper into session state for queries
                if faiss_store is not None:
                    st.session_state.vector_store = faiss_store
                    st.session_state.chunks = faiss_store.chunks
                    st.success(f"Processed {len(faiss_store.chunks)} chunks (including cached ones).")
                else:
                    st.warning("No new chunks were processed and no existing index found.")
    
    # -------------------------------
    # Q&A SECTION
    # -------------------------------
    
    # Initialize history if not present
    if "history" not in st.session_state:
        st.session_state.history = []
    
    st.subheader("Ask a Question")
    
    question = st.text_input("Enter your question about the documents:")
    
    # Ask button
    if st.button("Get Answer"):
        if question.strip():
            answer = answer_question_with_context(question, st.session_state.vector_store, st.session_state.chunks)
            if answer:
                # Save to history
                st.session_state.history.append({"q": question, "a": answer})
                st.success(answer)
            else:
                st.warning("No answer could be generated. Try rephrasing your question.")
        else:
            st.warning("Please enter a question.")
    
    # Display history
    if st.session_state.history:
        st.subheader("Conversation History")
        for i, item in enumerate(st.session_state.history, 1):
            st.markdown(f"**Q{i}:** {item['q']}")
            st.markdown(f"**A{i}:** {item['a']}")
            st.markdown("---")
    
    # Clear history button
    if st.button("Clear History"):
        st.session_state.history = []
        st.session_state.cleared = True  # flag for one rerun
    # Show success message only once
    if st.session_state.get("cleared", False):
        st.success("Conversation history cleared.")
        st.session_state.cleared = False
    
    st.markdown("---")
    st.write("This is test version by Thang Nguyen. Feedback is very much appreciated!")
